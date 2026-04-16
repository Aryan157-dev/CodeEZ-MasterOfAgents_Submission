import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from chart_generator import generate_chart, make_progress_rings
from image_fetcher import fetch_image, fetch_title_image, extract_keyword

# ── Brand colors ─────────────────────────────────────────────────────────
RED        = RGBColor(0xEF, 0x44, 0x44)
DARK       = RGBColor(0x2C, 0x2C, 0x2C)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GREY   = RGBColor(0xA0, 0xA0, 0xA0)
ACCENT     = RGBColor(0x1A, 0x1A, 0x2E)
CARD_BG    = RGBColor(0xF9, 0xF9, 0xF9)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

THEMES = {
    "red": {
        "primary": RGBColor(0xEF, 0x44, 0x44),
        "dark":    RGBColor(0x2C, 0x2C, 0x2C),
        "light":   RGBColor(0xF5, 0xF5, 0xF5),
        "mid":     RGBColor(0xA0, 0xA0, 0xA0),
    },
    "green": {
        "primary": RGBColor(0x4A, 0x7C, 0x59),
        "dark":    RGBColor(0x1E, 0x3A, 0x2F),
        "light":   RGBColor(0xF0, 0xF4, 0xF1),
        "mid":     RGBColor(0x7A, 0x9A, 0x84),
    },
    "blue": {
        "primary": RGBColor(0x1A, 0x6B, 0xB5),
        "dark":    RGBColor(0x0D, 0x2B, 0x4E),
        "light":   RGBColor(0xF0, 0xF4, 0xFA),
        "mid":     RGBColor(0x7A, 0x9A, 0xC0),
    },
}

def set_theme(theme_name="green"):
    global RED, DARK, LIGHT_GREY, MID_GREY
    t = THEMES.get(theme_name, THEMES["green"])
    RED        = t["primary"]
    DARK       = t["dark"]
    LIGHT_GREY = t["light"]
    MID_GREY   = t["mid"]


# ── Template layout indices ───────────────────────────────────────────────
# Accenture template layout map:
#   [0] 1_Cover         → title slide
#   [1] 2_Cover         → alternate cover
#   [2] Divider         → section divider
#   [3] Blank           → blank (use for all content slides)
#   [4] Title only      → title + content area
#   [5] Thank You       → conclusion/thank you
#
# UAE template layout map:
#   [0] 0_Title Company → title slide
#   [1] C_Section blue  → section divider
#   [2] 1_E_Title...    → title + subtitle + body
#   [3] 1_E_Title...    → body only
#   [4] 1_E_Title...    → blank with logo

LAYOUT_TITLE      = 0   # cover/title layout
LAYOUT_CONTENT    = 3   # blank layout — used for all content slides
LAYOUT_CONCLUSION = 5   # thank you layout (Accenture); _get_layout handles fallback


def _get_layout(prs, idx):
    """Safely get layout by index, fallback to blank."""
    layouts = prs.slide_layouts
    if idx < len(layouts):
        return layouts[idx]
    # fallback: find a blank layout
    for l in layouts:
        if "blank" in l.name.lower():
            return l
    return layouts[min(3, len(layouts) - 1)]


# ── Helpers ──────────────────────────────────────────────────────────────
def _add_text(slide, text, left, top, width, height,
              font_name="Calibri", font_size=14, bold=False,
              color=None, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name  = font_name
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color if color else DARK
    return txb


def _add_rect(slide, left, top, width, height, fill_color,
              line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def _add_circle(slide, left, top, size, fill_color, text, font_size=11):
    circle = slide.shapes.add_shape(9, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = fill_color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = str(text)
    run.font.name  = "Calibri"
    run.font.size  = Pt(font_size)
    run.font.bold  = True
    run.font.color.rgb = WHITE


def _slide_header(slide, title):
    """Minimal header that respects the template background — just title text + accent bar."""
    # Subtle dark band at top
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), DARK)
    # Left red strip
    _add_rect(slide, 0, 0, Inches(0.1), Inches(1.1), RED)
    # Title text
    _add_text(slide, title,
              left=Inches(0.3), top=Inches(0.15),
              width=Inches(11.0), height=Inches(0.8),
              font_size=24, bold=True, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════
# LAYOUT 1 — TITLE SLIDE
# ══════════════════════════════════════════════════════════════════════════
def build_title_slide(slide, title, subtitle, img=None):
    """Title slide — left side text, right side topic image."""
    # White base
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)

    # Right image panel
    if img:
        try:
            slide.shapes.add_picture(img,
                left=Inches(6.5), top=0,
                width=Inches(6.83), height=SLIDE_H)
        except Exception:
            _add_rect(slide, Inches(6.5), 0, Inches(6.83), SLIDE_H, DARK)
    else:
        _add_rect(slide, Inches(6.5), 0, Inches(6.83), SLIDE_H, DARK)

    # No overlay — let the image show through cleanly

    # Left red vertical strip
    _add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, RED)

    # Title
    _add_text(slide, title,
              left=Inches(0.4), top=Inches(1.8),
              width=Inches(5.8), height=Inches(2.2),
              font_name="Calibri", font_size=40, bold=True, color=RED)

    # Red divider line
    _add_rect(slide, Inches(0.4), Inches(4.15), Inches(4.0), Inches(0.06), RED)

    # Subtitle
    if subtitle:
        _add_text(slide, subtitle,
                  left=Inches(0.4), top=Inches(4.3),
                  width=Inches(5.8), height=Inches(1.0),
                  font_size=15, color=DARK)

    # Tag line
    _add_text(slide, "AI-Generated Presentation  •  Confidential",
              left=Inches(0.4), top=SLIDE_H - Inches(0.5),
              width=Inches(5.5), height=Inches(0.4),
              font_size=9, color=MID_GREY)


# ══════════════════════════════════════════════════════════════════════════
# LAYOUT 2 — EXECUTIVE SUMMARY
# ══════════════════════════════════════════════════════════════════════════
def build_executive_summary_slide(slide, title, key_points):
    _slide_header(slide, title)

    n       = min(len(key_points), 5)
    card_w  = (SLIDE_W - Inches(0.9)) / max(n, 1)
    card_h  = Inches(5.5)
    gap     = Inches(0.12)
    start_x = Inches(0.45)
    start_y = Inches(1.2)

    for i, point in enumerate(key_points[:n]):
        x = start_x + i * (card_w + gap / n)

        # Card shadow
        _add_rect(slide, x + Inches(0.04), start_y + Inches(0.04),
                  card_w - Inches(0.12), card_h, RGBColor(0xCC, 0xCC, 0xCC))
        # Card body
        _add_rect(slide, x, start_y, card_w - Inches(0.12), card_h, WHITE,
                  line_color=RGBColor(0xE8, 0xE8, 0xE8), line_width=Pt(1))
        # Top color bar
        bar_color = RED if i % 2 == 0 else DARK
        _add_rect(slide, x, start_y, card_w - Inches(0.12), Inches(0.1), bar_color)

        # Number circle
        _add_circle(slide, x + (card_w - Inches(0.12)) / 2 - Inches(0.28),
                    start_y + Inches(0.18), Inches(0.56), bar_color, str(i + 1), font_size=11)

        parts  = point.split(":", 1) if ":" in point else [point, ""]
        header = parts[0].strip()
        body   = parts[1].strip() if len(parts) > 1 else ""

        _add_text(slide, header,
                  left=x + Inches(0.1), top=start_y + Inches(0.9),
                  width=card_w - Inches(0.3), height=Inches(0.7),
                  font_size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)

        _add_rect(slide, x + Inches(0.3), start_y + Inches(1.68),
                  card_w - Inches(0.72), Inches(0.04), RGBColor(0xEE, 0xEE, 0xEE))

        if body:
            _add_text(slide, body,
                      left=x + Inches(0.12), top=start_y + Inches(1.82),
                      width=card_w - Inches(0.32), height=Inches(3.4),
                      font_size=10, color=MID_GREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# LAYOUT 3 — SPLIT PANEL
# ══════════════════════════════════════════════════════════════════════════
def _add_rect_transparent(slide, left, top, width, height, fill_color, alpha=128):
    """Add a rectangle with transparency (alpha 0=fully transparent, 255=opaque)."""
    from pptx.oxml.ns import qn
    from lxml import etree
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    # Set solid fill with alpha via XML
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    solidFill = spPr.find('.//' + qn('a:solidFill'))
    if solidFill is not None:
        spPr.remove(solidFill)
    # Build fill XML with alpha
    r, g, b = int(fill_color[0]), int(fill_color[1]), int(fill_color[2])
    alpha_pct = int((alpha / 255) * 100000)
    fill_xml = f'''<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <a:srgbClr val="{r:02X}{g:02X}{b:02X}">
        <a:alpha val="{alpha_pct}"/>
      </a:srgbClr>
    </a:solidFill>'''
    spPr.insert(0, etree.fromstring(fill_xml))
    return shape


def build_split_panel_slide(slide, title, key_points, context_text="", img=None):
    # Left dark panel base
    _add_rect(slide, 0, 0, Inches(4.2), SLIDE_H, DARK)

    # Add red strip
    _add_rect(slide, 0, 0, Inches(0.1), SLIDE_H, RED)

    # If image available, place it then add semi-transparent dark overlay
    if img:
        try:
            if hasattr(img, 'seek'):
                img.seek(0)
            slide.shapes.add_picture(img,
                left=Inches(0.1), top=0,
                width=Inches(4.1), height=SLIDE_H)
            # Semi-transparent overlay so text remains readable
            _add_rect_transparent(slide, Inches(0.1), 0, Inches(4.1), SLIDE_H,
                                   DARK, alpha=160)
        except Exception as e:
            print(f"   ⚠️  Could not add panel image: {e}")

    _add_text(slide, title,
              left=Inches(0.3), top=Inches(1.0),
              width=Inches(3.7), height=Inches(2.8),
              font_size=26, bold=True, color=WHITE)

    _add_rect(slide, Inches(0.3), Inches(4.0), Inches(2.0), Inches(0.06), RED)

    context = context_text or ""
    if context:
        _add_text(slide, context,
                  left=Inches(0.3), top=Inches(4.2),
                  width=Inches(3.7), height=Inches(2.8),
                  font_size=11, color=MID_GREY)

    right_x = Inches(4.5)
    right_w = SLIDE_W - right_x - Inches(0.25)

    n = min(len(key_points), 5)
    if n == 0:
        return

    available_h = SLIDE_H - Inches(0.3)
    gap    = Inches(0.12)
    card_h = (available_h - gap * (n - 1)) / n
    start_y = Inches(0.15)

    for i, point in enumerate(key_points[:n]):
        y = start_y + i * (card_h + gap)

        if ":" in point:
            parts = point.split(":", 1)
            label = parts[0].strip()
            desc  = parts[1].strip()
        else:
            label = ""
            desc  = point

        _add_rect(slide, right_x, y, right_w, card_h,
                  LIGHT_GREY if i % 2 == 0 else WHITE,
                  line_color=RGBColor(0xE8, 0xE8, 0xE8), line_width=Pt(0.5))
        _add_rect(slide, right_x, y, Inches(0.07), card_h, RED)
        _add_circle(slide, right_x + Inches(0.1), y + card_h / 2 - Inches(0.26),
                    Inches(0.52), RED, str(i + 1), font_size=9)

        if label:
            _add_text(slide, label,
                      left=right_x + Inches(0.75), top=y + Inches(0.08),
                      width=right_w - Inches(0.85), height=Inches(0.38),
                      font_size=13, bold=True, color=RED)
            _add_text(slide, desc,
                      left=right_x + Inches(0.75), top=y + Inches(0.42),
                      width=right_w - Inches(0.85), height=card_h - Inches(0.5),
                      font_size=11, color=DARK)
        else:
            _add_text(slide, desc,
                      left=right_x + Inches(0.75), top=y + Inches(0.12),
                      width=right_w - Inches(0.85), height=card_h - Inches(0.2),
                      font_size=12, color=DARK)


# ══════════════════════════════════════════════════════════════════════════
# LAYOUT 4 — TIMELINE (compact, sharp, no background)
# ══════════════════════════════════════════════════════════════════════════
def build_timeline_slide(slide, title, key_points):
    # Clean white background - no image
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    _slide_header(slide, title)

    n = min(len(key_points), 6)
    if n == 0:
        return

    # Tight layout constants
    start_x = Inches(0.4)
    total_w = SLIDE_W - Inches(0.8)
    item_w  = total_w / n
    spine_y = Inches(4.0)   # spine sits at vertical midpoint of content area

    # Draw spine line
    _add_rect(slide, start_x, spine_y, total_w, Inches(0.05), RGBColor(0xDD, 0xDD, 0xDD))
    # Red progress fill
    _add_rect(slide, start_x, spine_y, total_w, Inches(0.05), RED)

    for i, point in enumerate(key_points[:n]):
        parts = point.split(":", 1) if ":" in point else [f"Phase {i+1}", point]
        year  = parts[0].strip()
        event = parts[1].strip() if len(parts) > 1 else ""

        cx = start_x + i * item_w + item_w / 2

        # Node — white circle with red border
        ns = Inches(0.32)
        nx = cx - ns / 2
        ny = spine_y - ns / 2 + Inches(0.025)
        node = slide.shapes.add_shape(9, nx, ny, ns, ns)
        node.fill.solid()
        node.fill.fore_color.rgb = WHITE
        node.line.color.rgb = RED
        node.line.width = Pt(2.2)

        # Inner dot
        ds = Inches(0.13)
        dot = slide.shapes.add_shape(9, cx - ds/2, spine_y + Inches(0.025) - ds/2, ds, ds)
        dot.fill.solid()
        dot.fill.fore_color.rgb = RED
        dot.line.fill.background()

        tw = item_w - Inches(0.12)
        tx = cx - tw / 2

        if i % 2 == 0:
            # Year tag above spine
            _add_text(slide, year,
                      left=tx, top=Inches(1.35),
                      width=tw, height=Inches(0.42),
                      font_size=13, bold=True, color=RED,
                      align=PP_ALIGN.CENTER)
            # Thin connector
            _add_rect(slide, cx - Inches(0.01), Inches(1.77),
                      Inches(0.02), spine_y - Inches(1.77),
                      RGBColor(0xCC, 0xCC, 0xCC))
            # Card below spine
            cy, ch = spine_y + Inches(0.32), Inches(2.75)
            _add_rect(slide, tx, cy, tw, ch, LIGHT_GREY,
                      line_color=RGBColor(0xE0, 0xE0, 0xE0), line_width=Pt(0.7))
            _add_rect(slide, tx, cy, tw, Inches(0.05), RED)
            _add_text(slide, event,
                      left=tx + Inches(0.07), top=cy + Inches(0.12),
                      width=tw - Inches(0.14), height=ch - Inches(0.18),
                      font_size=9, color=DARK, align=PP_ALIGN.CENTER)
        else:
            # Card above spine
            ch = Inches(2.75)
            cy = spine_y - ch - Inches(0.32)
            _add_rect(slide, tx, cy, tw, ch, LIGHT_GREY,
                      line_color=RGBColor(0xE0, 0xE0, 0xE0), line_width=Pt(0.7))
            _add_rect(slide, tx, cy + ch - Inches(0.05), tw, Inches(0.05), RED)
            _add_text(slide, event,
                      left=tx + Inches(0.07), top=cy + Inches(0.08),
                      width=tw - Inches(0.14), height=ch - Inches(0.15),
                      font_size=9, color=DARK, align=PP_ALIGN.CENTER)
            # Thin connector
            _add_rect(slide, cx - Inches(0.01), cy + ch,
                      Inches(0.02), spine_y - cy - ch,
                      RGBColor(0xCC, 0xCC, 0xCC))
            # Year tag below spine
            _add_text(slide, year,
                      left=tx, top=spine_y + Inches(0.32),
                      width=tw, height=Inches(0.42),
                      font_size=13, bold=True, color=RED,
                      align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# LAYOUT 5 — KPI STATS
# ══════════════════════════════════════════════════════════════════════════
def build_kpi_stats_slide(slide, title, key_points):
    _slide_header(slide, title)

    n = min(len(key_points), 4)
    if n == 0:
        return

    card_w  = (SLIDE_W - Inches(1.0)) / n
    card_h  = Inches(5.5)
    start_x = Inches(0.5)
    start_y = Inches(1.2)

    for i, point in enumerate(key_points[:n]):
        parts   = [p.strip() for p in point.split(":", 2)]
        label   = parts[0] if len(parts) > 0 else ""
        value   = parts[1] if len(parts) > 1 else ""
        context = parts[2] if len(parts) > 2 else ""

        x = start_x + i * card_w

        # Shadow
        _add_rect(slide, x + Inches(0.05), start_y + Inches(0.05),
                  card_w - Inches(0.15), card_h, RGBColor(0xDD, 0xDD, 0xDD))
        # Card
        bg = LIGHT_GREY if i % 2 == 0 else WHITE
        _add_rect(slide, x, start_y, card_w - Inches(0.15), card_h, bg,
                  line_color=RGBColor(0xE0, 0xE0, 0xE0), line_width=Pt(1))
        # Top color bar
        _add_rect(slide, x, start_y, card_w - Inches(0.15), Inches(0.12), RED)

        _add_text(slide, label.upper(),
                  left=x + Inches(0.1), top=start_y + Inches(0.22),
                  width=card_w - Inches(0.3), height=Inches(0.5),
                  font_size=11, bold=True, color=MID_GREY, align=PP_ALIGN.CENTER)

        _add_rect(slide, x + Inches(0.3), start_y + Inches(0.8),
                  card_w - Inches(0.75), Inches(0.04), RGBColor(0xDD, 0xDD, 0xDD))

        _add_text(slide, value,
                  left=x + Inches(0.08), top=start_y + Inches(0.85),
                  width=card_w - Inches(0.2), height=Inches(1.8),
                  font_size=26, bold=True, color=RED, align=PP_ALIGN.CENTER)

        # Progress bar
        bar_w_total = card_w - Inches(0.6)
        _add_rect(slide, x + Inches(0.25), start_y + Inches(2.7),
                  bar_w_total, Inches(0.1), RGBColor(0xEE, 0xEE, 0xEE))
        fill_pct = max(0.1, 0.75 - (i * 0.1))
        _add_rect(slide, x + Inches(0.25), start_y + Inches(2.7),
                  bar_w_total * fill_pct, Inches(0.1), RED)

        if context:
            _add_text(slide, context,
                      left=x + Inches(0.1), top=start_y + Inches(2.95),
                      width=card_w - Inches(0.3), height=Inches(2.3),
                      font_size=11, color=DARK, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# LAYOUT 6 — TWO COLUMN COMPARE
# ══════════════════════════════════════════════════════════════════════════
def build_two_col_compare_slide(slide, title, key_points, metadata=None):
    if metadata is None:
        metadata = {}

    left_label  = metadata.get("left_label", "Column A")
    right_label = metadata.get("right_label", "Column B")

    _slide_header(slide, title)

    mid     = SLIDE_W / 2
    col_w   = mid - Inches(0.7)
    start_y = Inches(1.25)

    half = max(1, len(key_points) // 2)
    left_points  = key_points[:half]
    right_points = key_points[half:half * 2]

    for col_idx, (label, points, x_start) in enumerate([
        (left_label,  left_points,  Inches(0.5)),
        (right_label, right_points, mid + Inches(0.2))
    ]):
        col_color = RED if col_idx == 0 else DARK

        _add_rect(slide, x_start, start_y, col_w, Inches(0.55), col_color)
        _add_text(slide, label.upper(),
                  left=x_start + Inches(0.1), top=start_y + Inches(0.08),
                  width=col_w - Inches(0.2), height=Inches(0.4),
                  font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        item_h = Inches(0.88)
        gap    = Inches(0.12)
        for j, point in enumerate(points[:4]):
            y = start_y + Inches(0.65) + j * (item_h + gap)
            _add_rect(slide, x_start, y, col_w, item_h, LIGHT_GREY,
                      line_color=RGBColor(0xE0, 0xE0, 0xE0), line_width=Pt(0.5))
            _add_rect(slide, x_start, y, Inches(0.07), item_h, col_color)
            _add_text(slide, point,
                      left=x_start + Inches(0.18), top=y + Inches(0.1),
                      width=col_w - Inches(0.25), height=item_h - Inches(0.15),
                      font_size=11, color=DARK)

    _add_rect(slide, mid - Inches(0.05), start_y, Inches(0.06),
              SLIDE_H - start_y - Inches(0.3), RGBColor(0xE0, 0xE0, 0xE0))


# ══════════════════════════════════════════════════════════════════════════
# LAYOUT 7 — GRID 4 COLUMN
# ══════════════════════════════════════════════════════════════════════════
def build_grid_4col_slide(slide, title, key_points):
    _slide_header(slide, title)

    while len(key_points) < 4:
        key_points.append("—")

    col_w   = (SLIDE_W - Inches(0.8)) / 4
    col_h   = Inches(5.5)
    start_x = Inches(0.4)
    start_y = Inches(1.2)
    gap     = Inches(0.07)

    col_colors = [RED, DARK, RGBColor(0x55, 0x55, 0x55), RGBColor(0x77, 0x77, 0x77)]

    for i, point in enumerate(key_points[:4]):
        x = start_x + i * (col_w + gap)

        parts  = point.split(":", 1) if ":" in point else [f"Category {i+1}", point]
        header = parts[0].strip()
        body   = parts[1].strip() if len(parts) > 1 else ""

        _add_rect(slide, x, start_y, col_w, col_h, WHITE,
                  line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))
        _add_rect(slide, x, start_y, col_w, Inches(0.6), col_colors[i])
        _add_text(slide, header,
                  left=x + Inches(0.08), top=start_y + Inches(0.08),
                  width=col_w - Inches(0.16), height=Inches(0.48),
                  font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        items = [b.strip() for b in re.split(r'[|\n]', body) if b.strip()] if body else []
        if not items and body:
            items = [body]

        item_h = Inches(0.7)
        for j, item in enumerate(items[:5]):
            iy = start_y + Inches(0.7) + j * (item_h + Inches(0.05))
            _add_rect(slide, x + Inches(0.08), iy,
                      col_w - Inches(0.16), item_h, LIGHT_GREY,
                      line_color=RGBColor(0xEE, 0xEE, 0xEE), line_width=Pt(0.5))
            _add_text(slide, item,
                      left=x + Inches(0.14), top=iy + Inches(0.08),
                      width=col_w - Inches(0.28), height=item_h - Inches(0.1),
                      font_size=10, color=DARK, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# LAYOUT 8 — DATA TABLE
# ══════════════════════════════════════════════════════════════════════════
def build_data_table_slide(slide, title, table_data):
    _slide_header(slide, title)

    if not table_data or not table_data.get("headers"):
        _add_text(slide, "No table data available.",
                  left=Inches(0.5), top=Inches(2.0),
                  width=Inches(12), height=Inches(1.0),
                  font_size=14, color=MID_GREY)
        return

    headers = table_data["headers"]
    rows    = table_data["rows"][:9]  # max 9 rows
    n_cols  = len(headers)
    n_rows  = len(rows)

    # Dynamic sizing — fill the slide
    start_x   = Inches(0.4)
    start_y   = Inches(1.25)
    available = SLIDE_H - start_y - Inches(0.9)  # leave room for insight bar
    header_h  = Inches(0.65)
    row_h     = available / max(n_rows, 1)
    row_h     = min(row_h, Inches(0.85))  # cap at 0.85"
    col_w     = (SLIDE_W - Inches(0.8)) / n_cols

    # Detect numeric columns for highlighting
    def is_numeric(val):
        return bool(re.match(r'^[\d,.\+\-\$%xX]+$', str(val).strip()))

    # Header row
    for j, h in enumerate(headers):
        x = start_x + j * col_w
        # First col gets darker accent
        hcol = DARK if j == 0 else RED
        _add_rect(slide, x, start_y, col_w, header_h, hcol)
        _add_text(slide, str(h).upper(),
                  left=x + Inches(0.1), top=start_y + Inches(0.1),
                  width=col_w - Inches(0.15), height=header_h - Inches(0.1),
                  font_size=11, bold=True, color=WHITE)

    # Data rows
    for i, row in enumerate(rows):
        y      = start_y + header_h + i * row_h
        bg_col = LIGHT_GREY if i % 2 == 0 else WHITE
        for j, cell in enumerate(row[:n_cols]):
            x = start_x + j * col_w
            _add_rect(slide, x, y, col_w, row_h, bg_col,
                      line_color=RGBColor(0xE0, 0xE0, 0xE0), line_width=Pt(0.5))
            # Left border accent on first col
            if j == 0:
                _add_rect(slide, x, y, Inches(0.05), row_h, RED)
            cell_text = str(cell)
            font_bold  = (j == 0)
            font_color = RED if (j > 0 and is_numeric(cell)) else DARK
            _add_text(slide, cell_text,
                      left=x + Inches(0.12), top=y + Inches(0.08),
                      width=col_w - Inches(0.18), height=row_h - Inches(0.12),
                      font_size=11, bold=font_bold, color=font_color)

    # Insight bar at bottom
    bar_y = start_y + header_h + n_rows * row_h + Inches(0.1)
    if bar_y + Inches(0.5) < SLIDE_H:
        _add_rect(slide, start_x, bar_y, SLIDE_W - Inches(0.8), Inches(0.45), DARK)
        _add_rect(slide, start_x, bar_y, Inches(0.08), Inches(0.45), RED)
        insight = f"{n_rows} records  ·  {n_cols} metrics  ·  Source: Document Analysis"
        _add_text(slide, insight,
                  left=start_x + Inches(0.25), top=bar_y + Inches(0.08),
                  width=SLIDE_W - Inches(1.2), height=Inches(0.32),
                  font_size=10, color=MID_GREY)


# ══════════════════════════════════════════════════════════════════════════
# LAYOUT 9 — CHART
# ══════════════════════════════════════════════════════════════════════════
def build_chart_slide(slide, title, chart_path):
    """Chart slide: left info panel + chart on right."""
    # Left dark panel
    _add_rect(slide, 0, 0, Inches(3.0), SLIDE_H, DARK)
    _add_rect(slide, 0, 0, Inches(0.1), SLIDE_H, RED)

    _add_text(slide, title,
              left=Inches(0.25), top=Inches(1.0),
              width=Inches(2.6), height=Inches(2.5),
              font_size=20, bold=True, color=WHITE)

    _add_rect(slide, Inches(0.25), Inches(3.65), Inches(1.6), Inches(0.06), RED)

    _add_text(slide, "Data visualization\nfrom document analysis",
              left=Inches(0.25), top=Inches(3.85),
              width=Inches(2.6), height=Inches(1.0),
              font_size=9, color=MID_GREY)

    _add_circle(slide, Inches(0.85), Inches(5.5), Inches(0.8), RED, "📊", font_size=14)

    # Chart area
    chart_x = Inches(3.15)
    chart_w = SLIDE_W - chart_x - Inches(0.15)

    _add_rect(slide, chart_x, Inches(0.1), chart_w, SLIDE_H - Inches(0.2), LIGHT_GREY)

    if chart_path and os.path.exists(chart_path):
        slide.shapes.add_picture(
            chart_path,
            left=chart_x + Inches(0.1),
            top=Inches(0.2),
            width=chart_w - Inches(0.2),
            height=SLIDE_H - Inches(0.4)
        )
    else:
        _add_text(slide, "Chart could not be generated.",
                  left=chart_x, top=Inches(3.0),
                  width=chart_w, height=Inches(1.0),
                  font_size=14, color=MID_GREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# LAYOUT 10 — KPI VISUAL (progress rings)
# ══════════════════════════════════════════════════════════════════════════
def build_kpi_visual_slide(slide, title, parsed_doc):
    """Dark dashboard slide with circular progress rings."""
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK)
    _add_rect(slide, 0, 0, Inches(0.1), SLIDE_H, RED)
    _add_rect(slide, Inches(0.1), 0, SLIDE_W - Inches(0.1), Inches(1.1),
              RGBColor(0x1A, 0x1A, 0x1A))

    _add_text(slide, title,
              left=Inches(0.3), top=Inches(0.15),
              width=Inches(12.5), height=Inches(0.8),
              font_size=24, bold=True, color=WHITE)

    _add_rect(slide, Inches(0.3), Inches(1.12), Inches(3.5), Inches(0.06), RED)

    stats = parsed_doc.get("key_stats", [])
    rings_data = []
    for idx, s in enumerate(stats[:4]):
        raw  = str(s["value"])
        unit = "%" if "%" in raw else \
               "B" if "B" in raw or "billion" in raw.lower() else \
               "M" if "M" in raw or "million" in raw.lower() else ""
        nums = re.findall(r'[\d]+\.?\d*', raw)
        if nums:
            try:
                val = float(nums[0])
                fill_pcts = [0.85, 0.65, 0.75, 0.45]
                max_val   = val / fill_pcts[idx % 4]
                rings_data.append((s["label"][:12], val, max_val, unit))
            except:
                continue

    if not rings_data and parsed_doc.get("tables"):
        for table in parsed_doc["tables"][:3]:
            for row in table.get("rows", [])[:4]:
                for cell in row[1:]:
                    nums = re.findall(r'[\d]+\.?\d*', str(cell))
                    if nums:
                        try:
                            val = float(nums[0])
                            if val > 0:
                                rings_data.append((str(row[0])[:12], val, val * 1.5, ""))
                                break
                        except:
                            continue
                if len(rings_data) >= 4:
                    break
            if len(rings_data) >= 4:
                break

    if rings_data:
        chart_path = make_progress_rings(rings_data, "rings.png")
        if chart_path and os.path.exists(chart_path):
            slide.shapes.add_picture(
                chart_path,
                left=Inches(0.5), top=Inches(1.4),
                width=Inches(12.3), height=Inches(5.7)
            )
    else:
        _add_text(slide, "Key metrics visualization unavailable for this document.",
                  left=Inches(0.5), top=Inches(3.5),
                  width=Inches(12), height=Inches(1.0),
                  font_size=14, color=MID_GREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# LAYOUT 11 — CONCLUSION
# ══════════════════════════════════════════════════════════════════════════
def build_conclusion_slide(slide, title, key_points):
    """
    Uses the template's Thank You / conclusion layout background.
    We inject title + takeaway points on top.
    """
    # Semi-transparent dark overlay so our text is readable over template bg
    _add_rect(slide, 0, 0, Inches(6.5), SLIDE_H, RGBColor(0x1E, 0x1E, 0x1E))
    _add_rect(slide, 0, 0, Inches(0.1), SLIDE_H, RED)

    _add_text(slide, title,
              left=Inches(0.4), top=Inches(0.3),
              width=Inches(6.0), height=Inches(0.9),
              font_size=28, bold=True, color=WHITE)

    _add_rect(slide, Inches(0.4), Inches(1.25), Inches(3.2), Inches(0.06), RED)

    for i, point in enumerate(key_points[:5]):
        y = Inches(1.5) + i * Inches(1.0)
        _add_rect(slide, Inches(0.4), y, Inches(5.9), Inches(0.85),
                  RGBColor(0x3C, 0x3C, 0x3C))
        _add_rect(slide, Inches(0.4), y, Inches(0.08), Inches(0.85), RED)
        _add_text(slide, "→",
                  left=Inches(0.62), top=y + Inches(0.14),
                  width=Inches(0.5), height=Inches(0.5),
                  font_size=15, bold=True, color=RED)
        _add_text(slide, point,
                  left=Inches(1.2), top=y + Inches(0.14),
                  width=Inches(5.0), height=Inches(0.6),
                  font_size=13, color=WHITE)


# ── Helpers ───────────────────────────────────────────────────────────────
def add_slide_number(slide, number):
    _add_text(slide, str(number),
              left=SLIDE_W - Inches(0.65), top=SLIDE_H - Inches(0.42),
              width=Inches(0.5), height=Inches(0.35),
              font_size=9, color=MID_GREY, align=PP_ALIGN.RIGHT)


def find_table(parsed_doc, title_hint):
    if not parsed_doc.get("tables"):
        return None
    if not title_hint:
        return parsed_doc["tables"][0]

    title_hint_lower = title_hint.lower()

    for t in parsed_doc["tables"]:
        if t["title"].lower() == title_hint_lower:
            return t

    hint_words  = set(title_hint_lower.split())
    best_match  = None
    best_score  = 0
    for t in parsed_doc["tables"]:
        table_words = set(t["title"].lower().split())
        score = len(hint_words & table_words)
        if score > best_score:
            best_score = score
            best_match = t

    if best_score >= 2:
        return best_match

    return parsed_doc["tables"][0]


def _resolve_template(template_path, md_file=None):
    """
    Return a valid template path.
    - If caller supplied one and it exists → use it.
    - Otherwise match by markdown filename (word overlap).
    - Falls back to first Template_*.pptx found.
    - Returns None if nothing found.
    """
    if template_path and os.path.exists(template_path):
        return template_path

    templates = [f for f in os.listdir(".")
                 if f.lower().startswith("template_") and f.lower().endswith(".pptx")]

    if not templates:
        return None

    # Try to match by word overlap with the markdown filename
    if md_file:
        md_stem = os.path.splitext(os.path.basename(md_file))[0].lower()
        md_words = set(re.split(r'[\s_\-]+', md_stem))

        best_match = None
        best_score = 0
        for fname in templates:
            t_stem  = os.path.splitext(fname)[0].lower().replace("template_", "")
            t_words = set(re.split(r'[\s_\-]+', t_stem))
            score   = len(md_words & t_words)
            if score > best_score:
                best_score = score
                best_match = fname

        if best_match:
            print(f"   📁 Matched template: {best_match} (score={best_score})")
            return best_match

    # Fallback: just use the first one
    print(f"   📁 Using first available template: {templates[0]}")
    return templates[0]


# ══════════════════════════════════════════════════════════════════════════
# LAYOUT NEW-A — BIG STATEMENT
# ══════════════════════════════════════════════════════════════════════════
def build_big_statement_slide(slide, title, key_points):
    """
    Full-slide bold statement layout. Title = the big statement.
    key_points[0] = supporting insight line.
    key_points[1:] = 3 supporting fact pills at the bottom.
    """
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    # Left red strip
    _add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, RED)
    # Top accent bar
    _add_rect(slide, Inches(0.12), 0, SLIDE_W - Inches(0.12), Inches(0.12), DARK)
    # Bottom dark band
    _add_rect(slide, 0, SLIDE_H - Inches(1.2), SLIDE_W, Inches(1.2), DARK)

    # Big statement — vertically centered
    _add_text(slide, title,
              left=Inches(0.6), top=Inches(1.5),
              width=Inches(12.1), height=Inches(2.8),
              font_size=44, bold=True, color=DARK,
              align=PP_ALIGN.CENTER)

    # Red accent line under statement
    _add_rect(slide, Inches(4.0), Inches(4.3), Inches(5.3), Inches(0.08), RED)

    # Supporting insight
    if key_points:
        _add_text(slide, key_points[0],
                  left=Inches(0.6), top=Inches(4.5),
                  width=Inches(12.1), height=Inches(0.6),
                  font_size=16, color=MID_GREY,
                  align=PP_ALIGN.CENTER)

    # Fact pills along the bottom band
    facts = key_points[1:4] if len(key_points) > 1 else []
    if facts:
        pill_w = (SLIDE_W - Inches(1.0)) / len(facts)
        for i, fact in enumerate(facts):
            px = Inches(0.5) + i * pill_w
            py = SLIDE_H - Inches(1.05)
            _add_rect(slide, px + Inches(0.05), py,
                      pill_w - Inches(0.1), Inches(0.8),
                      RGBColor(0x2A, 0x2A, 0x2A))
            _add_rect(slide, px + Inches(0.05), py,
                      Inches(0.06), Inches(0.8), RED)
            _add_text(slide, fact,
                      left=px + Inches(0.2), top=py + Inches(0.12),
                      width=pill_w - Inches(0.35), height=Inches(0.6),
                      font_size=11, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════
# LAYOUT NEW-B — ICON ROW (numbered pillars)
# ══════════════════════════════════════════════════════════════════════════
def build_icon_row_slide(slide, title, key_points):
    """
    3-5 numbered pillars/steps displayed horizontally with large numbers.
    Great for processes, strategies, or key factors.
    """
    _slide_header(slide, title)

    n = min(len(key_points), 5)
    if n == 0:
        return

    card_w  = (SLIDE_W - Inches(0.6)) / n
    start_x = Inches(0.3)
    start_y = Inches(1.3)
    card_h  = SLIDE_H - start_y - Inches(0.2)

    col_colors = [RED, DARK, RGBColor(0x4A, 0x6A, 0x4A),
                  RGBColor(0x3A, 0x3A, 0x5A), RGBColor(0x5A, 0x3A, 0x3A)]

    for i, point in enumerate(key_points[:n]):
        x     = start_x + i * card_w
        color = col_colors[i % len(col_colors)]

        parts  = point.split(":", 1) if ":" in point else [f"Step {i+1}", point]
        header = parts[0].strip()
        body   = parts[1].strip() if len(parts) > 1 else ""

        # Card background
        _add_rect(slide, x + Inches(0.05), start_y, card_w - Inches(0.1), card_h,
                  WHITE, line_color=RGBColor(0xE0, 0xE0, 0xE0), line_width=Pt(1))

        # Top color band
        _add_rect(slide, x + Inches(0.05), start_y,
                  card_w - Inches(0.1), Inches(0.08), color)

        # Large number — smaller and tighter
        _add_text(slide, f"0{i+1}" if i < 9 else str(i+1),
                  left=x + Inches(0.1), top=start_y + Inches(0.18),
                  width=card_w - Inches(0.2), height=Inches(0.9),
                  font_size=38, bold=True, color=color,
                  align=PP_ALIGN.CENTER)

        # Divider
        _add_rect(slide, x + Inches(0.3), start_y + Inches(1.15),
                  card_w - Inches(0.7), Inches(0.04),
                  RGBColor(0xE0, 0xE0, 0xE0))

        # Header
        _add_text(slide, header,
                  left=x + Inches(0.1), top=start_y + Inches(1.28),
                  width=card_w - Inches(0.2), height=Inches(0.65),
                  font_size=13, bold=True, color=DARK,
                  align=PP_ALIGN.CENTER)

        # Body text — bigger font, more space
        if body:
            _add_text(slide, body,
                      left=x + Inches(0.12), top=start_y + Inches(2.05),
                      width=card_w - Inches(0.24), height=card_h - Inches(2.2),
                      font_size=12, color=MID_GREY,
                      align=PP_ALIGN.CENTER)


def _add_slide_background(slide, img):
    """Add a full-slide background image with dark overlay for readability."""
    if not img:
        return
    try:
        if hasattr(img, 'seek'):
            img.seek(0)
        # Insert picture at bottom of shape tree so it's behind everything
        pic = slide.shapes.add_picture(img, left=0, top=0,
                                        width=SLIDE_W, height=SLIDE_H)
        # Move picture to back by reordering XML
        sp_tree = slide.shapes._spTree
        sp_tree.remove(pic._element)
        sp_tree.insert(2, pic._element)  # index 2 = behind all other shapes
        # Add semi-transparent dark overlay using XML alpha
        _add_rect_transparent(slide, 0, 0, SLIDE_W, SLIDE_H,
                               RGBColor(0x0A, 0x0A, 0x0A), alpha=175)
    except Exception as e:
        print(f"   ⚠️  Background image failed: {e}")
def build_presentation(slide_plan, parsed_doc, template_path, output_path):
    resolved = _resolve_template(template_path)

    if resolved:
        print(f"   ✅ Using template: {resolved}")
        prs = Presentation(resolved)
        # Remove all existing slides from the template
        # (we keep the slide master / layouts but start fresh)
        xml_slides = prs.slides._sldIdLst
        for sld in list(xml_slides):
            xml_slides.remove(sld)
    else:
        print("   ⚠️  No template found — building with blank presentation")
        prs = Presentation()

    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print(f"\n🏗  Building {len(slide_plan)} slides...")

    chart_path = None
    doc_title  = parsed_doc.get("title", "")

    # Pre-fetch title image
    title_img = fetch_title_image(doc_title)

    for slide_info in slide_plan:
        layout_type = slide_info.get("layout_type", "split_panel")
        title       = slide_info.get("title", "")
        points      = slide_info.get("key_points", [])
        meta        = slide_info.get("metadata", {})
        use_tbl     = slide_info.get("use_table")
        num         = slide_info.get("slide_number", 1)

        print(f"   Slide {num}: [{layout_type}] {title[:50]}")

        # Pick the right template layout
        if layout_type == "title":
            layout = _get_layout(prs, LAYOUT_TITLE)
        elif layout_type == "conclusion":
            layout = _get_layout(prs, LAYOUT_CONCLUSION)
        else:
            layout = _get_layout(prs, LAYOUT_CONTENT)

        slide = prs.slides.add_slide(layout)

        # Remove ALL placeholder shapes from the slide so they don't show through
        sp_tree = slide.shapes._spTree
        for ph in slide.placeholders:
            try:
                sp_tree.remove(ph._element)
            except Exception:
                pass

        if layout_type == "title":
            subtitle = points[0] if points else parsed_doc.get("subtitle", "")
            build_title_slide(slide, title, subtitle, img=title_img)
        elif layout_type == "executive_summary":
            build_executive_summary_slide(slide, title, points)
        elif layout_type == "split_panel":
            context   = meta.get("context", "")
            kw        = extract_keyword(title, doc_title)
            slide_img = fetch_image(kw)
            build_split_panel_slide(slide, title, points, context, img=slide_img)
        elif layout_type == "timeline":
            build_timeline_slide(slide, title, points)
        elif layout_type == "kpi_stats":
            build_kpi_stats_slide(slide, title, points)
        elif layout_type == "two_col_compare":
            build_two_col_compare_slide(slide, title, points, meta)
        elif layout_type == "grid_4col":
            build_grid_4col_slide(slide, title, points)
        elif layout_type == "data_table":
            table_data = find_table(parsed_doc, use_tbl or title)
            build_data_table_slide(slide, title, table_data)
        elif layout_type == "chart":
            chart_filename = f"chart_{num}.png"
            table_data = find_table(parsed_doc, use_tbl or title)
            print(f"      DEBUG use_tbl: {use_tbl}")
            print(f"      DEBUG matched table: {table_data['title'] if table_data else None}")
            chart_path = generate_chart(parsed_doc, chart_filename, table_data)
            build_chart_slide(slide, title, chart_path)
        elif layout_type == "kpi_visual":
            build_kpi_visual_slide(slide, title, parsed_doc)
        elif layout_type == "big_statement":
            kw = extract_keyword(title, doc_title)
            bg = fetch_image(kw)
            _add_slide_background(slide, bg)
            build_big_statement_slide(slide, title, points)
        elif layout_type == "icon_row":
            build_icon_row_slide(slide, title, points)
        elif layout_type == "conclusion":
            kw = extract_keyword(title, doc_title)
            bg = fetch_image(kw)
            _add_slide_background(slide, bg)
            build_conclusion_slide(slide, title, points)
        else:
            kw        = extract_keyword(title, doc_title)
            slide_img = fetch_image(kw)
            build_split_panel_slide(slide, title, points, img=slide_img)

        if layout_type != "title":
            add_slide_number(slide, num)

    prs.save(output_path)
    print(f"\n✅ Saved: {output_path}")


# ── Entry point ───────────────────────────────────────────────────────────
if __name__ == "__main__":
    import sys
    from parser import parse_markdown
    from slide_planner import plan_slides

    API_KEY = os.environ.get("GROQ_API_KEY")

    MD_FILE = sys.argv[1] if len(sys.argv) > 1 else "Accenture Tech Acquisition Analysis.md"
    OUTPUT  = sys.argv[2] if len(sys.argv) > 2 else MD_FILE.replace(".md", ".pptx")

    print(f"Input:  {MD_FILE}")
    print(f"Output: {OUTPUT}")

    doc  = parse_markdown(MD_FILE)
    plan = plan_slides(doc, API_KEY)
    build_presentation(plan, doc, None, OUTPUT)